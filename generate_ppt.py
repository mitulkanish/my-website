from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def apply_slide_theme(slide):
    # Set a dark sleek background color by drawing a full-slide rectangle 
    # (more reliable than slide.background in pptx)
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    bg_shape = slide.shapes.add_shape(1, left, top, width, height) # 1 = msoShapeRectangle
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(15, 23, 42) # Deep slate/blue
    bg_shape.line.fill.background() # No border
    
    # Send it to back by modifying XML order.
    # We'll just rely on adding it FIRST in the script.
    
def format_title(title):
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(56, 189, 248) # Cyan accent
    title.text_frame.paragraphs[0].font.bold = True

def format_text(content):
    if not hasattr(content, 'paragraphs'): return
    for paragraph in content.paragraphs:
        if paragraph.font:
            paragraph.font.color.rgb = RGBColor(226, 232, 240) # Slate-200 Light text
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(226, 232, 240)

def create_presentation():
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_theme(slide)
    
    # We need to re-fetch or re-create text frames over the background shape
    # The simplest hack for a template is just leaving the text as-is and drawing the bg 
    # FIRST before applying text so it layers underneath. Since we already requested placeholders,
    # the placeholders naturally sit on top in PPT!
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "IntelliRisk: Academic Intelligence Platform"
    subtitle.text = "Empowering Education with Predictive Machine Learning"
    format_title(title)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184)
    
    # Slide 2: The Problem
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_theme(slide)
    title = slide.shapes.title
    title.text = "The Problem: Reactive Education"
    format_title(title)
    
    content = slide.placeholders[1].text_frame
    content.text = "Traditional Education is Reactive, Not Proactive"
    p = content.add_paragraph()
    p.text = "Teachers rely on delayed midterm and final exam results to identify struggling students."
    p = content.add_paragraph()
    p.text = "Attendance tracking is disconnected from actual academic performance data."
    p = content.add_paragraph()
    p.text = "Current software provides passive data (spreadsheets) rather than actionable insights."
    p = content.add_paragraph()
    p.text = "By the time an intervention occurs, students are already failing."
    format_text(content)

    # Slide 3: Our Solution
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_theme(slide)
    title = slide.shapes.title
    title.text = "Our Solution: Real-Time Tracking"
    format_title(title)
    
    content = slide.placeholders[1].text_frame
    content.text = "Real-Time Academic Health Tracking"
    p = content.add_paragraph()
    p.text = "IntelliRisk is a centralized platform connecting student attendance, test scores, and practical capabilities."
    p = content.add_paragraph()
    p.text = "We use Machine Learning to analyze data live and generate an 'Academic Health Score' out of 100."
    p = content.add_paragraph()
    p.text = "The system acts as an early warning radar, transforming raw classroom data into predictive success pipelines."
    format_text(content)

    # Slide 4: Key Feature
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_theme(slide)
    title = slide.shapes.title
    title.text = "Key Feature: Predictive AI"
    format_title(title)
    
    content = slide.placeholders[1].text_frame
    content.text = "Powered by Predictive AI"
    p = content.add_paragraph()
    p.text = "Our customized backend ML API analyzes 1,000+ data points instantly. It evaluates Mathematics, Computer Theory, Design Engineering, and C++ Programming alongside Theory & Practical Attendance."
    p = content.add_paragraph()
    p.text = "Classifies Profiles: Automatically categorizes students (e.g., 'Practical Learner', 'At-Risk')."
    p = content.add_paragraph()
    p.text = "Generates Insights: Tells students exactly what they are doing wrong and how to fix it."
    format_text(content)

    # Slide 5: Student Experience
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_theme(slide)
    title = slide.shapes.title
    title.text = "The Student Experience"
    format_title(title)
    
    content = slide.placeholders[1].text_frame
    content.text = "Personalized Learning Dashboards"
    p = content.add_paragraph()
    p.text = "Every student receives a beautiful, intuitive portal containing their unique dataset."
    p = content.add_paragraph()
    p.text = "Success Intelligence Tab: Features interactive Radar Charts breaking down practical vs. theoretical dimensions."
    p = content.add_paragraph()
    p.text = "A centralized hub for Weekly Tests, Quizzes, and personalized AI Alerts."
    format_text(content)
    
    if os.path.exists("dashboard_screenshot.png"):
        slide.shapes.add_picture("dashboard_screenshot.png", Inches(5), Inches(2), width=Inches(4.5))
    
    # Slide 6: Admin View
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_theme(slide)
    title = slide.shapes.title
    title.text = "The Administrator View"
    format_title(title)
    
    content = slide.placeholders[1].text_frame
    content.text = "God-Eye View for Faculty & Staff"
    p = content.add_paragraph()
    p.text = "A secure, dedicated portal strictly for administrators."
    p = content.add_paragraph()
    p.text = "Batch Analytics: Instantly visualize the entire classroom's health, averages, and profile distributions."
    p = content.add_paragraph()
    p.text = "At-Risk Ledger: Immediately highlights struggling students."
    format_text(content)
    
    if os.path.exists("dashboard_screenshot.png"):
        slide.shapes.add_picture("dashboard_screenshot.png", Inches(0.5), Inches(4), width=Inches(9))

    # Slide 7: Admin Deep Dive Navigation
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_theme(slide)
    title = slide.shapes.title
    title.text = "Admin Deep Dive Navigation"
    format_title(title)
    
    content = slide.placeholders[1].text_frame
    content.text = "Limitless Data Drill-Downs"
    p = content.add_paragraph()
    p.text = "Administrators aren't just given a summary—they have the power to click into any metric for a precise classroom breakdown."
    p = content.add_paragraph()
    p.text = "Class Directory: Search and filter the entire 1,000-student database."
    p = content.add_paragraph()
    p.text = "Individual Emulation: Click any name to instantly view that exact student’s Success Intelligence predictions."
    p = content.add_paragraph()
    p.text = "Subject Analysis Grid: Rank the entire class dynamically by subject performance."
    format_text(content)

    # Slide 8: One-Click Export
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_theme(slide)
    title = slide.shapes.title
    title.text = "One-Click Export & Reporting"
    format_title(title)
    
    content = slide.placeholders[1].text_frame
    content.text = "Automated Faculty Reporting"
    p = content.add_paragraph()
    p.text = "Say goodbye to manual data entry. IntelliRisk features a one-click 'Export Batch Report' engine."
    p = content.add_paragraph()
    p.text = "Generates instant offline CSV files."
    p = content.add_paragraph()
    p.text = "Compiles global metrics and profile distributions."
    p = content.add_paragraph()
    p.text = "Automatically lists the exact at-risk and excellent performing students for easy email distribution."
    format_text(content)

    # Slide 9: Technology Stack
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_theme(slide)
    title = slide.shapes.title
    title.text = "The Technology Stack"
    format_title(title)
    
    content = slide.placeholders[1].text_frame
    content.text = "Built for Speed and Scale"
    p = content.add_paragraph()
    p.text = "Frontend: React.js, Vite, and React-Router for a lightning-fast, glassmorphic SPA."
    p = content.add_paragraph()
    p.text = "UI/UX: Lucide Icons and Recharts for premium data visualization."
    p = content.add_paragraph()
    p.text = "Backend: Python and FastAPI ensuring rapid Machine Learning processing."
    p = content.add_paragraph()
    p.text = "Data Science: Scikit-Learn trained on extensive synthetic classroom datasets."
    format_text(content)

    # Slide 10: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    apply_slide_theme(slide)
    title = slide.shapes.title
    title.text = "Conclusion"
    format_title(title)
    
    content = slide.placeholders[1].text_frame
    content.text = "The Future of Intelligent Classrooms"
    p = content.add_paragraph()
    p.text = "Stop reacting to failures. Start predicting success."
    p = content.add_paragraph()
    p.text = "Thank you! Any questions?"
    format_text(content)

    prs.save('IntelliRisk_Presentation.pptx')
    print("Presentation generated successfully at IntelliRisk_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
